import fitz  # PyMuPDF
from pptx import Presentation
from app.models.schemas import ParsedPage, ParsedSlide, DocumentBlock, ImageAsset
from pathlib import Path

class IngestionService:
    def parse_pdf(self, file_path: Path) -> list[ParsedPage]:
        pages = []
        with fitz.open(file_path) as doc:
            for i, page in enumerate(doc):
                blocks = self._process_pdf_blocks(page)
                images = self._process_pdf_images(page, i)
                
                pages.append(ParsedPage(
                    page_number=i + 1,
                    blocks=blocks,
                    images=images
                ))
        return pages

    def _process_pdf_blocks(self, page) -> list[DocumentBlock]:
        blocks = []
        text_dict = page.get_text("dict")
        for b in text_dict["blocks"]:
            if b["type"] == 0:  # text
                for line in b["lines"]:
                    # Join spans without extra spaces to handle character-level fragmentation
                    raw_text = "".join(s["text"] for s in line["spans"])
                    # Normalize whitespace
                    clean_text = " ".join(raw_text.split()).strip()
                    
                    if not clean_text:
                        continue
                    
                    max_size = max(s["size"] for s in line["spans"])
                    kind = "heading" if max_size > 14 else "paragraph"
                    blocks.append(DocumentBlock(kind=kind, text=clean_text))
        return blocks

    def _process_pdf_images(self, page, page_index: int) -> list[ImageAsset]:
        images = []
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            images.append(ImageAsset(
                page_number=page_index + 1,
                index=img_index,
                width=img[2],
                height=img[3],
                extension=img[8]
            ))
        return images

    def parse_pptx(self, file_path: Path) -> list[ParsedSlide]:
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            title = slide.shapes.title.text if slide.shapes.title else None
            body = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    body.append(shape.text)
            
            speaker_notes = ""
            if slide.has_notes_slide:
                speaker_notes = slide.notes_slide.notes_text_frame.text
            
            slides.append(ParsedSlide(
                slide_number=i + 1,
                title=title,
                body=body,
                speaker_notes=speaker_notes
            ))
        return slides

ingestion_service = IngestionService()

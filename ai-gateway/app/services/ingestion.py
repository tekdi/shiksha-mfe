import fitz  # PyMuPDF
from pptx import Presentation
import io

class DocumentIngestionService:
    @staticmethod
    def extract_text_from_pdf(file_bytes: bytes) -> str:
        """
        Extracts text from a PDF file using PyMuPDF.
        """
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text_content = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text_content.append(page.get_text("text"))
        return "\n\n".join(text_content)

    @staticmethod
    def extract_text_from_ppt(file_bytes: bytes) -> str:
        """
        Extracts text and speaker notes from a PPTX file.
        """
        prs = Presentation(io.BytesIO(file_bytes))
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = [f"--- Slide {slide_num + 1} ---"]
            
            # Extract main text
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    slide_text.append(f"Speaker Notes: {notes}")
            
            text_content.append("\n".join(slide_text))
            
        return "\n\n".join(text_content)

    @staticmethod
    def process_document(filename: str, file_bytes: bytes) -> str:
        """
        Determines the file type and routes to the appropriate extractor.
        """
        if filename.lower().endswith('.pdf'):
            return DocumentIngestionService.extract_text_from_pdf(file_bytes)
        elif filename.lower().endswith(('.ppt', '.pptx')):
            return DocumentIngestionService.extract_text_from_ppt(file_bytes)
        else:
            raise ValueError(f"Unsupported file format for {filename}")

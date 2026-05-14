from pptx import Presentation
from typing import Dict, Any

def parse_pptx(file_path: str) -> Dict[str, Any]:
    """
    Parses a PPTX file to extract text and speaker notes per slide.
    """
    prs = Presentation(file_path)
    
    full_text = ""
    slides = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        
        # Extract Speaker Notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text
                
        slides.append({
            "slide_number": i + 1,
            "content": slide_text.strip(),
            "speaker_notes": notes_text.strip()
        })
        
        full_text += slide_text + "\n" + notes_text + "\n"
        
    return {
        "total_slides": len(slides),
        "slides": slides,
        "full_text": full_text.strip()
    }
